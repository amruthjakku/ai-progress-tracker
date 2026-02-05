"""
File Preview Service - Convert documents for web viewing
"""
import os
import io
from typing import Optional, Dict, Any


def get_file_preview(file_path: str, file_type: str, page: Optional[int] = None) -> bytes:
    """
    Generate preview content for different file types
    
    Args:
        file_path: Path to the file
        file_type: Extension (pdf, docx, pptx, ppt)
        page: Optional page number for paginated content
    
    Returns:
        bytes: Preview content
    """
    if file_type == "pdf":
        return _preview_pdf(file_path)
    elif file_type == "docx":
        return _preview_docx(file_path)
    elif file_type in ["pptx", "ppt"]:
        return _preview_pptx(file_path, page)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def get_preview_content_type(file_type: str) -> str:
    """Get the content type for preview response"""
    content_types = {
        "pdf": "application/pdf",
        "docx": "text/html",
        "pptx": "image/png",
        "ppt": "image/png"
    }
    return content_types.get(file_type, "application/octet-stream")


def get_file_info(file_path: str, file_type: str) -> Dict[str, Any]:
    """Get file metadata for preview rendering"""
    info = {
        "size_bytes": os.path.getsize(file_path) if os.path.exists(file_path) else 0
    }
    
    if file_type == "pdf":
        info["page_count"] = _get_pdf_page_count(file_path)
    elif file_type in ["pptx", "ppt"]:
        info["slide_count"] = _get_pptx_slide_count(file_path)
    
    return info


# ============ PDF Preview ============
def _preview_pdf(file_path: str) -> bytes:
    """Return PDF file directly for iframe embedding"""
    with open(file_path, "rb") as f:
        return f.read()


def _get_pdf_page_count(file_path: str) -> int:
    """Get number of pages in PDF"""
    try:
        # Use PyPDF2 if available for page count
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        return len(reader.pages)
    except ImportError:
        # Fallback: return -1 if PyPDF2 not installed
        return -1
    except Exception:
        return -1


# ============ DOCX Preview ============
def _preview_docx(file_path: str) -> bytes:
    """Convert DOCX to HTML for web display"""
    try:
        from docx import Document
        from docx.shared import Inches
        
        doc = Document(file_path)
        
        # Build HTML
        html_parts = [
            "<!DOCTYPE html>",
            "<html><head>",
            "<meta charset='utf-8'>",
            "<style>",
            "body { font-family: Arial, sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; line-height: 1.6; }",
            "h1, h2, h3 { color: #333; }",
            "p { margin: 10px 0; }",
            "table { border-collapse: collapse; width: 100%; margin: 15px 0; }",
            "td, th { border: 1px solid #ddd; padding: 8px; text-align: left; }",
            "th { background-color: #f4f4f4; }",
            "</style>",
            "</head><body>"
        ]
        
        for para in doc.paragraphs:
            style = para.style.name.lower() if para.style else ""
            text = para.text.strip()
            
            if not text:
                continue
            
            if "heading 1" in style:
                html_parts.append(f"<h1>{text}</h1>")
            elif "heading 2" in style:
                html_parts.append(f"<h2>{text}</h2>")
            elif "heading 3" in style:
                html_parts.append(f"<h3>{text}</h3>")
            else:
                html_parts.append(f"<p>{text}</p>")
        
        # Handle tables
        for table in doc.tables:
            html_parts.append("<table>")
            for i, row in enumerate(table.rows):
                html_parts.append("<tr>")
                for cell in row.cells:
                    tag = "th" if i == 0 else "td"
                    html_parts.append(f"<{tag}>{cell.text}</{tag}>")
                html_parts.append("</tr>")
            html_parts.append("</table>")
        
        html_parts.append("</body></html>")
        
        return "\n".join(html_parts).encode("utf-8")
    
    except ImportError:
        return b"<html><body><p>python-docx not installed. Cannot preview DOCX files.</p></body></html>"
    except Exception as e:
        return f"<html><body><p>Error previewing document: {str(e)}</p></body></html>".encode("utf-8")


# ============ PPTX Preview ============
def _preview_pptx(file_path: str, page: Optional[int] = None) -> bytes:
    """Convert PPTX slide to image"""
    try:
        from pptx import Presentation
        from PIL import Image
        
        prs = Presentation(file_path)
        slide_idx = (page or 1) - 1
        
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            slide_idx = 0
        
        # Create a simple representation of the slide
        slide = prs.slides[slide_idx]
        
        # Create an image representing the slide
        width, height = 1280, 720
        img = Image.new('RGB', (width, height), color='white')
        
        # For a proper implementation, we'd need to render shapes
        # This is a simplified version that creates a placeholder
        from PIL import ImageDraw, ImageFont
        draw = ImageDraw.Draw(img)
        
        # Extract text from slide
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        
        # Draw text on image
        y_position = 50
        for text in texts[:10]:  # Limit to first 10 text elements
            # Truncate long text
            display_text = text[:100] + "..." if len(text) > 100 else text
            draw.text((50, y_position), display_text, fill='black')
            y_position += 40
        
        # Add slide number
        draw.text((width - 100, height - 50), f"Slide {slide_idx + 1}/{len(prs.slides)}", fill='gray')
        
        # Convert to bytes
        img_buffer = io.BytesIO()
        img.save(img_buffer, format='PNG')
        img_buffer.seek(0)
        
        return img_buffer.read()
    
    except ImportError:
        # Return a placeholder image
        from PIL import Image, ImageDraw
        img = Image.new('RGB', (800, 600), color='#f0f0f0')
        draw = ImageDraw.Draw(img)
        draw.text((200, 280), "python-pptx not installed", fill='red')
        
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        buffer.seek(0)
        return buffer.read()
    except Exception as e:
        from PIL import Image, ImageDraw
        img = Image.new('RGB', (800, 600), color='#fff0f0')
        draw = ImageDraw.Draw(img)
        draw.text((50, 280), f"Error: {str(e)[:60]}", fill='red')
        
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        buffer.seek(0)
        return buffer.read()


def _get_pptx_slide_count(file_path: str) -> int:
    """Get number of slides in PPTX"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        return len(prs.slides)
    except Exception:
        return -1
